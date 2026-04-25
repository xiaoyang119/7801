"""
Generate academic PPTX presentation from MRM Validation Report.
Output: outputs/MRM_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Color palette (academic / minimal) ────────────────────────────────────────
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x00, 0x00, 0x00)
DGRAY    = RGBColor(0x33, 0x33, 0x33)
MGRAY    = RGBColor(0x59, 0x59, 0x59)
LGRAY    = RGBColor(0xF2, 0xF2, 0xF2)
HGRAY    = RGBColor(0xD9, 0xD9, 0xD9)   # table header fill
ALTROW   = RGBColor(0xF7, 0xF7, 0xF7)
ACCENT   = RGBColor(0x1A, 0x30, 0x55)   # dark navy — used sparingly
RULE     = RGBColor(0x88, 0x88, 0x88)   # thin horizontal rule color
PASS_F   = RGBColor(0xD6, 0xEA, 0xD6)
PASS_T   = RGBColor(0x1A, 0x5C, 0x1A)
AMBER_F  = RGBColor(0xFF, 0xF2, 0xCC)
AMBER_T  = RGBColor(0x7B, 0x4F, 0x00)
ISSUE_F  = RGBColor(0xFC, 0xE4, 0xD6)
ISSUE_T  = RGBColor(0x8B, 0x1A, 0x00)

PROJ_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_DIR  = os.path.join(PROJ_DIR, "outputs")
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=Pt(12), bold=False,
             color=DGRAY, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox


def slide_header(slide, title, subtitle=None):
    """Clean academic header: dark accent bar on left edge + title."""
    # thin left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT)
    # thin top rule
    add_rect(slide, Inches(0.08), Inches(0), SLIDE_W - Inches(0.08), Inches(0.06), HGRAY)
    add_text(slide, title,
             Inches(0.25), Inches(0.12), Inches(12.8), Inches(0.65),
             font_size=Pt(28), bold=True, color=ACCENT)
    if subtitle:
        add_rect(slide, Inches(0.08), Inches(0.78), SLIDE_W - Inches(0.08), Inches(0.03), HGRAY)
        add_text(slide, subtitle,
                 Inches(0.25), Inches(0.82), Inches(12.5), Inches(0.32),
                 font_size=Pt(12), color=MGRAY, italic=True)


def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)


def table_on_slide(slide, headers, rows, l, t, w, h,
                   header_font=Pt(10), row_font=Pt(9.5)):
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols, l, t, w, h).table

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HGRAY
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(hdr)
        run.font.bold = True
        run.font.size = header_font
        run.font.color.rgb = DGRAY
        run.font.name = "Calibri"

    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else ALTROW
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = tf.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size = row_font
            run.font.color.rgb = DGRAY
            run.font.name = "Calibri"

    return tbl


def status_cell(tbl, ri, ci, status):
    fills = {"Pass": (PASS_F, PASS_T), "Review": (AMBER_F, AMBER_T), "Issue": (ISSUE_F, ISSUE_T)}
    if status in fills:
        f, t = fills[status]
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = f
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = t
        run.font.bold = True


def slide_footer(slide, text="FRE 7801 — Model Risk Management  |  NYU Tandon  |  Spring 2026"):
    add_rect(slide, Inches(0.08), Inches(7.28), SLIDE_W - Inches(0.08), Inches(0.02), HGRAY)
    add_text(slide, text, Inches(0.25), Inches(7.3), Inches(12.8), Inches(0.18),
             font_size=Pt(8), color=MGRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # left accent bar
    add_rect(s, Inches(0), Inches(0), Inches(0.1), SLIDE_H, ACCENT)
    # thin horizontal divider
    add_rect(s, Inches(0.1), Inches(3.5), SLIDE_W - Inches(0.1), Inches(0.02), HGRAY)

    add_text(s, "Model Risk Assessment of Machine Learning-Based Credit Scoring",
             Inches(0.4), Inches(0.8), Inches(12.6), Inches(1.4),
             font_size=Pt(34), bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

    add_text(s, "A Case Study Using the Give Me Some Credit Dataset",
             Inches(0.4), Inches(2.25), Inches(12), Inches(0.55),
             font_size=Pt(18), color=MGRAY, align=PP_ALIGN.LEFT, italic=True)

    add_text(s, "Model Validation Report",
             Inches(0.4), Inches(2.9), Inches(12), Inches(0.48),
             font_size=Pt(14), color=MGRAY, align=PP_ALIGN.LEFT)

    add_text(s,
             "Xiaoyang Zhang (xz5476)  ·  Yifan Xu (yx3598)  ·  Yike Ma (ym2552)",
             Inches(0.4), Inches(3.7), Inches(12), Inches(0.42),
             font_size=Pt(13), color=DGRAY)

    add_text(s, "FRE 7801 — Model Risk Management  |  NYU Tandon School of Engineering",
             Inches(0.4), Inches(4.18), Inches(12), Inches(0.38),
             font_size=Pt(12), color=MGRAY)

    add_text(s, "Spring 2026  |  Validation Standard: SR 11-7",
             Inches(0.4), Inches(4.6), Inches(12), Inches(0.38),
             font_size=Pt(12), color=MGRAY)

    # risk tier label (understated)
    add_rect(s, Inches(0.4), Inches(5.5), Inches(3.2), Inches(0.5), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Risk Tier: Tier 1 — High Risk  (MRS 88/100)",
             Inches(0.5), Inches(5.52), Inches(3.0), Inches(0.44),
             font_size=Pt(11), bold=True, color=ISSUE_T)


def slide_agenda(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Agenda", "Model Risk Assessment — LightGBM Credit Scoring")

    items = [
        ("01", "Executive Summary"),
        ("02", "Dataset & Scope"),
        ("03", "Data Quality Review"),
        ("04", "Model Architecture"),
        ("05", "Quantitative Validation"),
        ("06", "Overfitting Analysis"),
        ("07", "Sensitivity Analysis"),
        ("08", "Fair Lending Analysis"),
        ("09", "Model Risk Tiering & SR 11-7 Checklist"),
        ("10", "Final Validation Opinion"),
    ]

    for i, (num, label) in enumerate(items):
        col = i % 2
        row_i = i // 2
        lx = Inches(0.5) + col * Inches(6.5)
        ty = Inches(1.3) + row_i * Inches(1.0)

        add_rect(s, lx, ty + Inches(0.05), Inches(0.48), Inches(0.42), HGRAY,
                 line_color=RULE, line_width=Pt(0.5))
        add_text(s, num, lx, ty + Inches(0.05), Inches(0.48), Inches(0.42),
                 font_size=Pt(11), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, label, lx + Inches(0.58), ty + Inches(0.08),
                 Inches(5.7), Inches(0.42), font_size=Pt(14), color=DGRAY)

    slide_footer(s)


def slide_exec_summary(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Executive Summary", "Key Validation Findings")

    headers = ["Validation Area", "Result", "Assessment"]
    rows = [
        ("Champion Validation AUC",        "0.8601",  "Strong discriminatory power"),
        ("Challenger Validation AUC",       "0.8317",  "Strong benchmark"),
        ("Champion KS Statistic",          "0.5700",  "Excellent rank ordering"),
        ("Champion Train AUC",             "0.9245",  "High in-sample fit"),
        ("Champion AUC Gap (single split)", "0.0644", "Review finding"),
        ("CV Mean AUC Gap (5-fold)",        "0.1669", "Overfitting confirmed"),
        ("Max PSI",                         "0.0005", "Stable on random split"),
        ("Young-Borrower DIR (LGBM)",       "0.3351", "ECOA-permitted; reflects 4.3x default rate diff."),
        ("Internal MRS Score",              "88/100", "Tier 1 — High Risk"),
    ]

    tbl = table_on_slide(s, headers, rows,
                         Inches(0.25), Inches(1.2), Inches(8.6), Inches(5.75),
                         row_font=Pt(10.5))

    # conclusion box
    add_rect(s, Inches(9.1), Inches(1.2), Inches(4.0), Inches(2.5), LGRAY,
             line_color=ACCENT, line_width=Pt(1))
    add_text(s, "Validation Conclusion",
             Inches(9.2), Inches(1.25), Inches(3.8), Inches(0.38),
             font_size=Pt(12), bold=True, color=ACCENT)
    add_text(s,
             "Conditionally approved pending overfitting remediation and governance controls.\n\n"
             "Age use is ECOA-compliant.",
             Inches(9.2), Inches(1.65), Inches(3.8), Inches(1.0),
             font_size=Pt(11), color=DGRAY)

    add_rect(s, Inches(9.1), Inches(3.9), Inches(4.0), Inches(1.35), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Model Under Review",
             Inches(9.2), Inches(3.95), Inches(3.8), Inches(0.38),
             font_size=Pt(11), bold=True, color=MGRAY)
    add_text(s,
             "LightGBM Credit Scoring Pipeline\nGive Me Some Credit — Kaggle 2011",
             Inches(9.2), Inches(4.35), Inches(3.8), Inches(0.85),
             font_size=Pt(10.5), color=MGRAY)

    add_rect(s, Inches(9.1), Inches(5.4), Inches(4.0), Inches(1.1), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Validation Standard",
             Inches(9.2), Inches(5.45), Inches(3.8), Inches(0.38),
             font_size=Pt(11), bold=True, color=MGRAY)
    add_text(s,
             "SR 11-7 — Federal Reserve & OCC",
             Inches(9.2), Inches(5.85), Inches(3.8), Inches(0.6),
             font_size=Pt(10.5), color=MGRAY)

    slide_footer(s)


def slide_dataset_scope(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Dataset & Scope", "Give Me Some Credit — Kaggle 2011")

    stats = [
        ("150,000", "Training Observations"),
        ("101,503", "Test Observations"),
        ("~6.7%",   "Default Rate"),
        ("10",      "Input Features"),
    ]
    for i, (val, label) in enumerate(stats):
        lx = Inches(0.25) + i * Inches(3.1)
        add_rect(s, lx, Inches(1.2), Inches(2.85), Inches(0.95), LGRAY,
                 line_color=RULE, line_width=Pt(0.5))
        add_text(s, val, lx, Inches(1.22), Inches(2.85), Inches(0.55),
                 font_size=Pt(26), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, label, lx, Inches(1.74), Inches(2.85), Inches(0.38),
                 font_size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)

    headers = ["Feature", "Risk Consideration"]
    rows = [
        ("RevolvingUtilizationOfUnsecuredLines", "Extreme outliers require caps"),
        ("age",                                  "Fair lending sensitivity"),
        ("NumberOfTime30-59DaysPastDueNotWorse", "Sentinel values require treatment"),
        ("DebtRatio",                            "Heavy right tail"),
        ("MonthlyIncome",                        "~20% missingness"),
        ("NumberOfOpenCreditLinesAndLoans",      "Nonlinear effect"),
        ("NumberOfTimes90DaysLate",              "Strong risk signal"),
        ("NumberRealEstateLoansOrLines",         "Ambiguous directional effect"),
        ("NumberOfTime60-89DaysPastDueNotWorse", "Sentinel values require treatment"),
        ("NumberOfDependents",                   "Missingness and weak sensitivity"),
    ]
    table_on_slide(s, headers, rows,
                   Inches(0.25), Inches(2.35), Inches(7.6), Inches(4.75),
                   row_font=Pt(9.5))

    add_rect(s, Inches(8.1), Inches(2.35), Inches(5.0), Inches(4.75), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Out of Scope", Inches(8.2), Inches(2.4), Inches(4.8), Inches(0.4),
             font_size=Pt(12), bold=True, color=MGRAY)
    oos = [
        "Legal fair lending determination",
        "Independent second-line validation",
        "Internal audit review",
        "Production implementation testing",
        "Live monitoring vs. bank portfolio",
        "Adverse action reason-code validation",
        "True out-of-time back-testing",
    ]
    for i, item in enumerate(oos):
        add_text(s, f"–  {item}", Inches(8.25), Inches(2.9) + i * Inches(0.3),
                 Inches(4.7), Inches(0.3), font_size=Pt(10.5), color=DGRAY)

    slide_footer(s)


def slide_data_quality(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Data Quality Review", "Missingness, Outliers & Population Stability")

    add_rect(s, Inches(0.25), Inches(1.2), Inches(4.6), Inches(1.1), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Missing Value Summary",
             Inches(0.35), Inches(1.23), Inches(4.4), Inches(0.36),
             font_size=Pt(11), bold=True, color=MGRAY)
    add_text(s,
             "MonthlyIncome: 19.82% missing (29,731 rows)\n"
             "NumberOfDependents: 2.62% missing (3,924 rows)\n"
             "Missingness indicators preserved as features",
             Inches(0.35), Inches(1.6), Inches(4.4), Inches(0.65),
             font_size=Pt(10.5), color=DGRAY)

    headers = ["Feature", "PSI", "Status"]
    rows = [
        ("DebtRatio",                            "0.0005", "Stable"),
        ("MonthlyIncome",                        "0.0004", "Stable"),
        ("RevolvingUtilization",                 "0.0003", "Stable"),
        ("age",                                  "0.0003", "Stable"),
        ("OpenCreditLines",                      "0.0002", "Stable"),
        ("RealEstateLoans",                      "0.0001", "Stable"),
        ("30-59DaysPastDue",                     "0.0000", "Stable"),
        ("90DaysLate",                           "0.0000", "Stable"),
        ("NumberOfDependents",                   "0.0000", "Stable"),
        ("MonthlyIncome_missing",                "0.0000", "Stable"),
    ]
    table_on_slide(s, headers, rows,
                   Inches(0.25), Inches(2.5), Inches(4.6), Inches(4.55),
                   row_font=Pt(9))

    add_image(s, os.path.join(OUT_DIR, "psi_chart.png"),
              Inches(5.1), Inches(1.2), Inches(8.0), Inches(5.9))

    add_text(s, "Max PSI 0.0005 — random split; not production drift evidence",
             Inches(0.35), Inches(7.05), Inches(4.5), Inches(0.22),
             font_size=Pt(8.5), color=MGRAY, italic=True)

    slide_footer(s)


def slide_model_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Model Architecture", "Champion vs. Challenger Pipeline")

    add_text(s, "Champion — LightGBM Pipeline",
             Inches(0.25), Inches(1.18), Inches(5.9), Inches(0.38),
             font_size=Pt(13), bold=True, color=ACCENT)
    add_rect(s, Inches(0.25), Inches(1.55), Inches(5.9), Inches(0.02), ACCENT)

    champion_steps = [
        ("1", "Missing Value Treatment",      "Impute + missingness indicator flags"),
        ("2", "Yeo-Johnson PowerTransformer", "Reduce skew in numerical features"),
        ("3", "MiniBatchKMeans Clustering",   "Borrower-level segment assignment"),
        ("4", "Weight of Evidence Encoding",  "WOE transform of cluster labels"),
        ("5", "Polynomial Features (deg-2)",  "Capture nonlinear interactions"),
        ("6", "SelectPercentile (ANOVA F)",   "Dimensionality reduction"),
        ("7", "LightGBM Classifier",          "Gradient-boosted tree w/ class weighting"),
    ]
    for i, (num, step, detail) in enumerate(champion_steps):
        ty = Inches(1.7) + i * Inches(0.72)
        add_rect(s, Inches(0.25), ty, Inches(0.4), Inches(0.55), HGRAY,
                 line_color=RULE, line_width=Pt(0.5))
        add_text(s, num, Inches(0.25), ty, Inches(0.4), Inches(0.55),
                 font_size=Pt(11), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, step, Inches(0.75), ty + Inches(0.02), Inches(5.3), Inches(0.28),
                 font_size=Pt(11), bold=True, color=DGRAY)
        add_text(s, detail, Inches(0.75), ty + Inches(0.3), Inches(5.3), Inches(0.25),
                 font_size=Pt(9.5), color=MGRAY, italic=True)

    add_text(s, "Challenger — Logistic Regression",
             Inches(6.5), Inches(1.18), Inches(6.6), Inches(0.38),
             font_size=Pt(13), bold=True, color=ACCENT)
    add_rect(s, Inches(6.5), Inches(1.55), Inches(6.6), Inches(0.02), ACCENT)

    add_rect(s, Inches(6.5), Inches(1.7), Inches(6.6), Inches(2.2), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    lr_text = (
        "–  Standardized inputs via StandardScaler\n"
        "–  Class weighting to handle imbalance\n"
        "–  Traditional scorecard-style benchmark\n"
        "–  Transparent, auditable, regulator-friendly\n"
        "–  Fully interpretable coefficients"
    )
    add_text(s, lr_text, Inches(6.65), Inches(1.8), Inches(6.3), Inches(2.0),
             font_size=Pt(12), color=DGRAY)

    add_text(s, "Champion — Key Risk Considerations",
             Inches(6.5), Inches(4.1), Inches(6.6), Inches(0.38),
             font_size=Pt(12), bold=True, color=MGRAY)
    add_rect(s, Inches(6.5), Inches(4.46), Inches(6.6), Inches(0.02), HGRAY)

    risks = [
        ("KMeans-WOE",   "Non-monotonic, difficult to audit"),
        ("Polynomial",   "Overfitting and interpretability risk"),
        ("LightGBM",     "Black-box model risk"),
        ("Missingness",  "May encode demographic differences"),
    ]
    for i, (component, risk) in enumerate(risks):
        ty = Inches(4.55) + i * Inches(0.56)
        add_text(s, component, Inches(6.6), ty, Inches(2.0), Inches(0.48),
                 font_size=Pt(11), bold=True, color=DGRAY)
        add_text(s, risk, Inches(8.8), ty, Inches(4.2), Inches(0.48),
                 font_size=Pt(11), color=MGRAY)

    slide_footer(s)


def slide_quant_results(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Quantitative Validation", "Benchmark Performance — Champion vs. Challenger")

    headers = ["Model", "AUC-ROC", "KS", "Gini", "PR-AUC", "F-beta", "Brier", "Train AUC", "AUC Gap"]
    rows = [
        ("LR Baseline",   "0.8317", "0.5075", "0.6635", "0.3344", "0.4619", "0.1652", "0.8291", "-0.0026"),
        ("LGBM Champion", "0.8601", "0.5700", "0.7202", "0.3927", "0.5147", "0.1253", "0.9245",  "0.0644"),
    ]

    tbl = table_on_slide(s, headers, rows,
                         Inches(0.25), Inches(1.15), Inches(12.8), Inches(1.35),
                         row_font=Pt(10.5), header_font=Pt(10.5))

    # flag the LGBM AUC Gap cell
    tbl.cell(2, 8).fill.solid()
    tbl.cell(2, 8).fill.fore_color.rgb = AMBER_F
    tbl.cell(2, 8).text_frame.paragraphs[0].runs[0].font.color.rgb = AMBER_T

    callouts = [
        ("+0.0284", "AUC Lift"),
        ("+0.0625", "KS Lift"),
        ("-0.0399", "Brier Reduction"),
        ("+0.0583", "PR-AUC Lift"),
    ]
    for i, (val, label) in enumerate(callouts):
        lx = Inches(0.25) + i * Inches(3.15)
        add_rect(s, lx, Inches(2.7), Inches(2.9), Inches(0.8), LGRAY,
                 line_color=RULE, line_width=Pt(0.5))
        add_text(s, val, lx, Inches(2.72), Inches(2.9), Inches(0.44),
                 font_size=Pt(20), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, label, lx, Inches(3.14), Inches(2.9), Inches(0.32),
                 font_size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)

    add_image(s, os.path.join(OUT_DIR, "roc_pr_curves.png"),
              Inches(0.25), Inches(3.65), Inches(8.1), Inches(3.45))

    add_rect(s, Inches(8.55), Inches(3.65), Inches(4.55), Inches(3.45), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Note on Threshold Selection",
             Inches(8.65), Inches(3.7), Inches(4.35), Inches(0.4),
             font_size=Pt(12), bold=True, color=MGRAY)
    add_text(s,
             "Champion is a strong rank-ordering model.\n\n"
             "Production threshold ≠ 0.50 mechanically.\n\n"
             "Threshold should reflect:\n"
             "–  Credit policy\n"
             "–  Expected loss\n"
             "–  Approval capacity\n"
             "–  Calibration",
             Inches(8.65), Inches(4.15), Inches(4.35), Inches(2.85),
             font_size=Pt(11), color=DGRAY)

    slide_footer(s)


def slide_overfitting(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Overfitting Analysis", "5-Fold Cross-Validation Evidence")

    for i, (name, train, val, gap, flag) in enumerate([
        ("LR Baseline",   "0.8291", "0.8317", "-0.0026", False),
        ("LGBM Champion", "0.9245", "0.8601", "+0.0644", True),
    ]):
        lx = Inches(0.25) + i * Inches(6.4)
        add_text(s, name, lx, Inches(1.15), Inches(6.1), Inches(0.38),
                 font_size=Pt(13), bold=True, color=ACCENT)
        add_rect(s, lx, Inches(1.5), Inches(6.1), Inches(0.02), HGRAY)

        for j, (label, metric) in enumerate([("Train AUC", train), ("Val AUC", val), ("AUC Gap", gap)]):
            tx = lx + j * Inches(2.0)
            gap_bg = AMBER_F if flag and label == "AUC Gap" else LGRAY
            gap_tc = AMBER_T if flag and label == "AUC Gap" else DGRAY
            add_rect(s, tx, Inches(1.6), Inches(1.9), Inches(0.85),
                     gap_bg, line_color=RULE, line_width=Pt(0.5))
            add_text(s, metric, tx, Inches(1.62), Inches(1.9), Inches(0.5),
                     font_size=Pt(20), bold=True, color=gap_tc, align=PP_ALIGN.CENTER)
            add_text(s, label, tx, Inches(2.1), Inches(1.9), Inches(0.3),
                     font_size=Pt(9), color=MGRAY, align=PP_ALIGN.CENTER)

        comment = "No overfitting — gap negligible" if not flag else "Exceeds 0.03 review threshold"
        add_text(s, comment, lx, Inches(2.58), Inches(6.1), Inches(0.35),
                 font_size=Pt(11), color=ISSUE_T if flag else PASS_T, italic=True)

    # CV results table
    add_text(s, "5-Fold CV Results — LGBM Champion",
             Inches(0.25), Inches(3.1), Inches(7.0), Inches(0.38),
             font_size=Pt(12), bold=True, color=MGRAY)
    cv_headers = ["Fold", "Train AUC", "Val AUC", "AUC Gap", "KS"]
    cv_rows = [
        ("Fold 1", "0.9968", "0.8358", "0.1610", "0.5325"),
        ("Fold 2", "0.9966", "0.8309", "0.1657", "0.5131"),
        ("Fold 3", "0.9969", "0.8186", "0.1783", "0.4925"),
        ("Fold 4", "0.9964", "0.8324", "0.1640", "0.5204"),
        ("Fold 5", "0.9962", "0.8308", "0.1654", "0.5123"),
        ("Mean",   "0.9966", "0.8297", "0.1669", "0.5142"),
    ]
    tbl = table_on_slide(s, cv_headers, cv_rows,
                         Inches(0.25), Inches(3.55), Inches(6.9), Inches(3.6),
                         row_font=Pt(10))
    # mean row highlight (row index 6 = 7th row, 0-based)
    for ci in range(5):
        tbl.cell(6, ci).fill.solid()
        tbl.cell(6, ci).fill.fore_color.rgb = HGRAY

    add_image(s, os.path.join(OUT_DIR, "cv_auc_folds.png"),
              Inches(7.4), Inches(3.1), Inches(5.7), Inches(4.1))

    slide_footer(s)


def slide_sensitivity(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Sensitivity Analysis", "±10% Input Perturbation — Average Absolute Score Movement")

    headers = ["Rank", "Feature", "Avg Abs. Score Movement"]
    rows = [
        ("1",  "age",                                  "0.013805"),
        ("2",  "RevolvingUtilizationOfUnsecuredLines",  "0.011690"),
        ("3",  "MonthlyIncome",                         "0.004820"),
        ("4",  "NumberOfOpenCreditLinesAndLoans",       "0.004510"),
        ("5",  "DebtRatio",                             "0.003750"),
        ("6",  "NumberOfTime60-89DaysPastDueNotWorse",  "0.001460"),
        ("7",  "NumberOfTimes90DaysLate",               "0.001335"),
        ("8",  "NumberRealEstateLoansOrLines",          "0.000965"),
        ("9",  "NumberOfDependents",                    "0.000225"),
        ("10", "NumberOfTime30-59DaysPastDueNotWorse",  "0.000175"),
    ]
    tbl = table_on_slide(s, headers, rows,
                         Inches(0.25), Inches(1.15), Inches(7.6), Inches(5.8),
                         row_font=Pt(10.5))

    # subtle highlight for age row
    for ci in range(3):
        tbl.cell(1, ci).fill.solid()
        tbl.cell(1, ci).fill.fore_color.rgb = AMBER_F

    add_rect(s, Inches(8.1), Inches(1.15), Inches(5.0), Inches(5.8), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Key Observation",
             Inches(8.2), Inches(1.2), Inches(4.8), Inches(0.4),
             font_size=Pt(13), bold=True, color=MGRAY)
    add_rect(s, Inches(8.2), Inches(1.6), Inches(4.8), Inches(0.02), HGRAY)
    add_text(s,
             "Age is the most sensitive feature.\n\n"
             "This is assessed in the context of fair lending "
             "in the following slide.\n\n"
             "Age is an ECOA-permitted factor in empirically "
             "derived credit scoring systems (15 U.S.C. § 1691).\n\n"
             "The sensitivity finding means age is doing genuine "
             "predictive work proportional to actual default rate "
             "differences across age groups.",
             Inches(8.2), Inches(1.7), Inches(4.8), Inches(5.0),
             font_size=Pt(12), color=DGRAY)

    slide_footer(s)


def slide_fair_lending(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Fair Lending Analysis", "Age-Based Disparate Impact — ECOA Context")

    add_rect(s, Inches(0.25), Inches(1.15), Inches(12.8), Inches(0.52), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s,
             "Age is an ECOA-permitted factor (15 U.S.C. § 1691) in empirically derived credit scoring systems. "
             "DIR gap is consistent with a 4.3× actual default rate difference across age groups.",
             Inches(0.35), Inches(1.18), Inches(12.6), Inches(0.46),
             font_size=Pt(10.5), color=DGRAY, italic=True)

    headers = ["Age Group", "Model", "N", "Default Rate", "Approval Rate", "AUC", "DIR vs. Best"]
    rows = [
        ("Young (<30)",    "LGBM Champion", "2,623",  "12.54%", "24.70%", "0.8189", "0.3351"),
        ("Young (<30)",    "LR Baseline",   "2,623",  "12.54%", "22.80%", "0.7814", "0.3051"),
        ("Middle (30-60)", "LGBM Champion", "27,795", "8.12%",  "39.94%", "0.8459", "0.5418"),
        ("Middle (30-60)", "LR Baseline",   "27,795", "8.12%",  "39.59%", "0.8173", "0.5297"),
        ("Senior (>60)",   "LGBM Champion", "14,582", "2.89%",  "73.72%", "0.8534", "1.0000"),
        ("Senior (>60)",   "LR Baseline",   "14,582", "2.89%",  "74.74%", "0.8131", "1.0000"),
    ]
    tbl = table_on_slide(s, headers, rows,
                         Inches(0.25), Inches(1.82), Inches(7.1), Inches(3.6),
                         row_font=Pt(10))

    add_image(s, os.path.join(OUT_DIR, "fairness_age_groups.png"),
              Inches(7.55), Inches(1.15), Inches(5.55), Inches(4.3))

    # LDA result box
    add_rect(s, Inches(0.25), Inches(5.6), Inches(7.1), Inches(1.55), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s, "Less-Discriminatory Alternative Test",
             Inches(0.35), Inches(5.65), Inches(6.9), Inches(0.38),
             font_size=Pt(11), bold=True, color=MGRAY)
    add_text(s,
             "Age-blind LGBM: DIR improves 0.3351 → 0.5081 (+0.173) but AUC drops −0.0372 "
             "(below LR baseline). Remaining DIR gap driven by age-correlated legitimate credit factors. "
             "Age use confirmed ECOA-compliant.",
             Inches(0.35), Inches(6.05), Inches(6.9), Inches(0.9),
             font_size=Pt(10.5), color=DGRAY)

    slide_footer(s)


def slide_risk_tiering(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Model Risk Tiering & SR 11-7 Checklist", "Internal Model Risk Scorecard")

    headers = ["Dimension", "Max", "Score", "%", "Justification"]
    rows = [
        ("Materiality",         "30", "27", "90%", "Credit decisions for large borrower population"),
        ("Regulatory Exposure", "30", "25", "83%", "SR 11-7, ECOA, explainability exposure"),
        ("Operational Factor",  "10",  "8", "80%", "Automated pipeline, no override control"),
        ("Model Complexity",    "30", "28", "93%", "LightGBM + WOE + polynomial features"),
        ("Total",              "100", "88", "88%", "Tier 1 — High Risk"),
    ]
    tbl = table_on_slide(s, headers, rows,
                         Inches(0.25), Inches(1.15), Inches(6.5), Inches(2.8),
                         row_font=Pt(10))
    # total row
    for ci in range(5):
        tbl.cell(5, ci).fill.solid()
        tbl.cell(5, ci).fill.fore_color.rgb = HGRAY

    add_image(s, os.path.join(OUT_DIR, "mrs_scorecard_chart.png"),
              Inches(0.25), Inches(4.1), Inches(6.5), Inches(3.1))

    # SR 11-7 checklist
    add_text(s, "SR 11-7 Validation Checklist",
             Inches(7.0), Inches(1.15), Inches(6.1), Inches(0.38),
             font_size=Pt(13), bold=True, color=MGRAY)
    add_rect(s, Inches(7.0), Inches(1.5), Inches(6.1), Inches(0.02), HGRAY)

    sr_headers = ["SR 11-7 Area", "Status", "Key Action"]
    sr_rows = [
        ("Purpose & Scope",    "Pass",   "Maintain model-use limits"),
        ("Data Quality",       "Pass",   "Document Kaggle data limitations"),
        ("Benchmarking",       "Pass",   "Explain complexity-performance tradeoff"),
        ("Outcomes Analysis",  "Review", "Simplify pipeline; apply early stopping"),
        ("Ongoing Monitoring", "Pass",   "Assign owner and monitoring cadence"),
        ("Fair Lending",       "Pass",   "Monitor DIR semiannually"),
        ("Governance",         "Pass",   "Document lack of independent validation"),
    ]
    sr_tbl = table_on_slide(s, sr_headers, sr_rows,
                            Inches(7.0), Inches(1.6), Inches(6.1), Inches(4.0),
                            row_font=Pt(10))
    for ri, row in enumerate(sr_rows, start=1):
        status_cell(sr_tbl, ri, 1, row[1])

    add_image(s, os.path.join(OUT_DIR, "sr117_risk_chart.png"),
              Inches(7.0), Inches(5.8), Inches(6.1), Inches(1.6))

    slide_footer(s)


def slide_final_opinion(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Final Validation Opinion", "")

    # Main statement
    add_rect(s, Inches(0.25), Inches(1.15), Inches(12.8), Inches(1.3), LGRAY,
             line_color=ACCENT, line_width=Pt(1))
    add_text(s, "Conditional Approval Pending Remediation",
             Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.55),
             font_size=Pt(22), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s,
             "Age use is ECOA-compliant. SHAP explanations, 5-fold CV, and LDA test completed. "
             "Remaining gaps: overfitting, governance controls, threshold definition.",
             Inches(0.4), Inches(1.72), Inches(12.5), Inches(0.65),
             font_size=Pt(12), color=DGRAY, align=PP_ALIGN.CENTER)

    items = [
        ("Strengths",
         ["– Strong discrimination: AUC 0.8601, KS 0.5700",
          "– Outperforms LR challenger on all metrics",
          "– ECOA-compliant age use documented",
          "– SHAP global and local explanations generated",
          "– Imputation assumptions empirically validated"]),
        ("Remaining Gaps",
         ["– CV mean AUC gap 0.1669 (overfitting)",
          "– No production score threshold defined",
          "– No override / human review policy",
          "– No independent validation performed",
          "– No live monitoring framework deployed"]),
        ("Academic Use",
         ["– Suitable for case-study demonstration",
          "– All SR 11-7 checkpoints addressed",
          "– Tier 1 — High Risk (MRS 88/100)",
          "– Requires independent validation before",
          "   any real-world lending application"]),
    ]
    for i, (title, bullets) in enumerate(items):
        lx = Inches(0.25) + i * Inches(4.3)
        add_rect(s, lx, Inches(2.65), Inches(4.1), Inches(0.42), HGRAY)
        add_text(s, title, lx, Inches(2.67), Inches(4.1), Inches(0.38),
                 font_size=Pt(12), bold=True, color=DGRAY, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            add_text(s, b, lx + Inches(0.1), Inches(3.18) + j * Inches(0.45),
                     Inches(4.0), Inches(0.42), font_size=Pt(10.5), color=DGRAY)

    add_rect(s, Inches(0.25), Inches(6.45), Inches(12.8), Inches(0.55), LGRAY,
             line_color=RULE, line_width=Pt(0.5))
    add_text(s,
             "Thank You  —  Questions Welcome",
             Inches(0.25), Inches(6.48), Inches(12.8), Inches(0.48),
             font_size=Pt(13), color=MGRAY, align=PP_ALIGN.CENTER)

    slide_footer(s)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)
    slide_agenda(prs)
    slide_exec_summary(prs)
    slide_dataset_scope(prs)
    slide_data_quality(prs)
    slide_model_architecture(prs)
    slide_quant_results(prs)
    slide_overfitting(prs)
    slide_sensitivity(prs)
    slide_fair_lending(prs)
    slide_risk_tiering(prs)
    slide_final_opinion(prs)

    out_path = os.path.join(OUT_DIR, "MRM_Presentation.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
